#!/usr/bin/env python3
"""Assemble the Virtue market-overview deck on the Indeed corporate template."""
import os, copy
import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
TMPL = "/Users/garyh/Downloads/Indeed Corporate Deck Template-Feb2026-POWERPOINT 2/Indeed Corporate Deck Template-Feb2026.pptx"
CH = os.path.join(HERE, "charts")
AS = os.path.join(HERE, "assets")
OUT = os.path.join(HERE, "Virtue Integrated Care - Market Overview.pptx")

# ---- Brand palette ----
BLUE   = RGBColor(0x00, 0x3A, 0x9B)
NAVY   = RGBColor(0x00, 0x12, 0x3C)
DEEP   = RGBColor(0x00, 0x1E, 0x57)
A11Y   = RGBColor(0x55, 0x92, 0xFD)
MAGENTA= RGBColor(0x9D, 0x2B, 0x6B)
PURPLE = RGBColor(0x56, 0x44, 0xBF)
TEAL   = RGBColor(0x23, 0x7E, 0xA3)
ORANGE = RGBColor(0xB4, 0x60, 0x2B)
DARK   = RGBColor(0x2D, 0x2D, 0x2D)
GREY   = RGBColor(0x6A, 0x71, 0x80)
LGREY  = RGBColor(0xE7, 0xEA, 0xF0)
PALE   = RGBColor(0xF4, 0xF6, 0xFA)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
FONT   = "Indeed Sans"
FONT_L = "Indeed Sans Light"

prs = Presentation(TMPL)
EMU_W, EMU_H = prs.slide_width, prs.slide_height

# ---- strip demo slides, keep masters/layouts/theme ----
# Remove both the sldId entries AND their relationships so the orphaned
# slide parts are not re-serialised (avoids duplicate part-name collisions).
xml_slides = prs.slides._sldIdLst
pres_part = prs.part
for sid in list(xml_slides):
    rId = sid.get(qn('r:id'))
    xml_slides.remove(sid)
    try:
        pres_part.drop_rel(rId)
    except Exception:
        pass
BLANK = prs.slide_masters[0].slide_layouts[0]   # 'Blank'

def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, l, t, w, h, fill=None, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w or 1)
    return sp

def txt(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=4, line_spacing=1.04, wrap=True):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    if isinstance(runs[0], tuple):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (text, size, color, bold, *rest) in para:
            r = p.add_run(); r.text = text
            r.font.size = Pt(size); r.font.bold = bold
            r.font.name = rest[0] if rest else FONT
            r.font.color.rgb = color
    return tb

def img(s, path, l, t, w=None, h=None):
    kw = {}
    if w: kw['width'] = Inches(w)
    if h: kw['height'] = Inches(h)
    return s.shapes.add_picture(path, Inches(l), Inches(t), **kw)

def img_fit(s, path, l, t, maxw, maxh, halign='center', valign='middle'):
    """Place image scaled to fit within box preserving aspect ratio."""
    from PIL import Image
    iw, ih = Image.open(path).size
    ar = iw / ih
    w = maxw; h = w / ar
    if h > maxh:
        h = maxh; w = h * ar
    lx = l + (maxw - w) * (0.5 if halign == 'center' else (1 if halign == 'right' else 0))
    ty = t + (maxh - h) * (0.5 if valign == 'middle' else (1 if valign == 'bottom' else 0))
    return s.shapes.add_picture(path, Inches(lx), Inches(ty), width=Inches(w), height=Inches(h))

# ----------------------------------------------------------------------
# Editable native-shape "charts" (so competitor names are editable in
# Google Slides as ordinary text boxes, not baked into an image).
# ----------------------------------------------------------------------
_wb = openpyxl.load_workbook(os.path.join(HERE, 'Market Data (1) copy.xlsx'), data_only=True)
def _load(name):
    ws = _wb[name]; h = [ws.cell(1, c).value for c in range(1, ws.max_column + 1)]
    out = []
    for r in range(2, ws.max_row + 1):
        d = {h[c - 1]: ws.cell(r, c).value for c in range(1, ws.max_column + 1)}
        if d['companyid']:
            out.append(d)
    return out
def _num(v):
    try: return float(v)
    except: return 0.0
def _short(c):
    return c.rsplit('[', 1)[0].strip()
_cur = _load('Jan to May 2026')
_virt = [r for r in _cur if 'virtue' in r['companyid'].lower()][0]
_tops = sorted(_cur, key=lambda r: _num(r['spend_local']), reverse=True)[:12]
_cd = sorted(_tops + [_virt], key=lambda r: _num(r['spend_local']))   # ascending -> Virtue at bottom
C_ITEMS = [(_short(r['companyid']), _num(r['spend_local'])/1000, 'virtue' in r['companyid'].lower()) for r in _cd]
_mkt = sum(_num(r['Tot_Applystarts']) for r in _cur) / sum(_num(r['Clicks']) for r in _cur) * 100
D_ITEMS = [('Virtue', _num(_virt['Applystart_rate'])*100, 'v')]
for _k, _lab in [('Comfort Keepers', 'Comfort Keepers'), ('Mowlam', 'Mowlam'), ('Dovida', 'Dovida'),
                 ('Be Independent', 'Be Independent'), ('Danu Home Care', 'Danu Home Care')]:
    for r in _cur:
        if _k.lower() in r['companyid'].lower():
            D_ITEMS.append((_lab, _num(r['Applystart_rate'])*100, 'c')); break
D_ITEMS.append(('Market avg', _mkt, 'm'))

def hbar_native(s, x, y, w, h, items, namew, labw, vfmt):
    n = len(items); rowh = h / n
    barh = min(rowh*0.6, 0.2)
    vmax = max(it[1] for it in items)
    bx = x + namew; bwmax = w - namew - labw
    rect(s, bx, y, 0.012, h, fill=LGREY)                      # value axis line
    for i, (name, val, hl) in enumerate(items):
        cyc = y + h - (i + 0.5)*rowh
        col = BLUE if hl else GREY
        txt(s, x, cyc-rowh/2, namew-0.08, rowh,
            [[(name, 8, (BLUE if hl else DARK), bool(hl))]],
            align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.9, space_after=0)
        bw = max(bwmax*val/vmax, 0.03)
        rect(s, bx, cyc-barh/2, bw, barh, fill=col)
        txt(s, bx+bw+0.06, cyc-rowh/2, labw-0.04, rowh,
            [[(vfmt(val), 7.6, (BLUE if hl else GREY), bool(hl))]],
            anchor=MSO_ANCHOR.MIDDLE, space_after=0)

def col_native(s, x, y, w, h, items, vfmt):
    n = len(items); cs = w / n; cw = cs*0.5
    vmax = max(it[1] for it in items)*1.16
    by = y + h
    rect(s, x, by, w, 0.012, fill=LGREY)                     # category axis line
    for i, (name, val, kind) in enumerate(items):
        cx = x + (i + 0.5)*cs
        bh = h*val/vmax
        col = BLUE if kind == 'v' else (NAVY if kind == 'm' else GREY)
        rect(s, cx-cw/2, by-bh, cw, bh, fill=col)
        txt(s, cx-cs/2, by-bh-0.26, cs, 0.24, [[(vfmt(val), 9, (BLUE if kind == 'v' else DARK), kind == 'v')]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM, space_after=0)
        txt(s, cx-cs/2, by+0.06, cs, 0.5, [[(name, 8.2, (BLUE if kind == 'v' else DARK), kind == 'v')]],
            align=PP_ALIGN.CENTER, line_spacing=0.92, space_after=0)

def notes(s, talk, perspective, takeaway, question=None):
    tf = s.notes_slide.notes_text_frame
    tf.text = ""
    def block(title, body):
        p = tf.add_paragraph(); 
        r = p.add_run(); r.text = title; r.font.bold = True; r.font.size = Pt(11)
        p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = body; r2.font.size = Pt(11)
        tf.add_paragraph()
    # reuse first empty paragraph
    block("Talk track:", talk)
    block("Perspective:", perspective)
    block("Key takeaway:", takeaway)
    if question:
        block("Client question:", question)

M = 0.55          # side margin
CW = 10 - 2*M     # content width

def content_header(s, eyebrow, headline, headline_size=21):
    rect(s, M, 0.52, 0.34, 0.07, fill=BLUE)             # accent tick
    txt(s, M, 0.30, CW, 0.3, [[(eyebrow.upper(), 10.5, BLUE, True)]])
    txt(s, M, 0.62, CW, 0.95, [[(headline, headline_size, DARK, True)]], line_spacing=1.02)

def footer(s, source, page):
    img(s, os.path.join(AS, "indeed_blue.png"), M, 5.27, w=0.62)
    if source:
        txt(s, 2.0, 5.30, 6.0, 0.25, [[("Source: " + source, 7.3, GREY, False)]])
    txt(s, 9.0, 5.30, 0.9, 0.25, [[(str(page), 8.5, GREY, False)]], align=PP_ALIGN.RIGHT)

def divider(num, title, sub, talk=None):
    s = slide()
    rect(s, -0.1, -0.1, 10.2, 5.825, fill=NAVY)
    # subtle accent block
    rect(s, 0, 0, 0.16, 5.625, fill=A11Y)
    txt(s, M+0.15, 1.45, 6, 1.4, [[(num, 64, A11Y, True, FONT_L)]])
    txt(s, M+0.15, 2.75, 8.5, 1.2, [[(title, 30, WHITE, True)]], line_spacing=1.0)
    if sub:
        txt(s, M+0.15, 3.85, 7.8, 0.9, [[(sub, 13, RGBColor(0xC7,0xD3,0xEC), False)]], line_spacing=1.15)
    img(s, os.path.join(AS, "indeed_white.png"), M+0.15, 4.85, w=0.7)
    if talk:
        s.notes_slide.notes_text_frame.text = "Talk track: " + talk
    return s

PAGE = [1]
def pg():
    PAGE[0] += 1; return PAGE[0]

# ======================================================================
# 1. TITLE
# ======================================================================
s = slide()
rect(s, -0.1, -0.1, 10.2, 5.825, fill=NAVY)
rect(s, 0, 0, 0.16, 5.625, fill=A11Y)
img(s, os.path.join(AS, "indeed_white.png"), M, 0.5, w=0.95)
txt(s, M, 1.95, 8.9, 0.4, [[("MARKET OVERVIEW  |  TALENT STRATEGY", 12, A11Y, True)]])
txt(s, M, 2.4, 8.9, 1.5, [[("Virtue Integrated Care Services", 38, WHITE, True)]], line_spacing=1.0)
txt(s, M, 3.55, 8.6, 0.9,
    [[("Turning a maternity and health-insurance benefit into a measurable", 15, RGBColor(0xD7,0xE0,0xF2), False)],
     [("recruitment and retention advantage in the Irish private care market", 15, RGBColor(0xD7,0xE0,0xF2), False)]],
    line_spacing=1.15)
txt(s, M, 5.0, 8.9, 0.35, [[("Prepared for the C-suite executive strategy session  ·  Confidential", 10.5, RGBColor(0x9F,0xB2,0xD6), False)]])
notes(s,
 "This deck gives Virtue's leadership a full view of the Irish private care recruitment market and shows where the new maternity top-up and health insurance benefit fits commercially. We'll move from market context, to where Virtue stands today, to where the candidates are, and finally to how we turn the benefit into a measurable advantage.",
 "Frame this as a commercial conversation, not an HR update. The benefit is an investment and this deck is about proving the return.",
 "Virtue has made a genuinely differentiating move - the job now is to activate it and measure it before competitors catch up.")

# ======================================================================
# 2. EXECUTIVE SUMMARY
# ======================================================================
s = slide()
content_header(s, "Executive summary", "Virtue has created a market-differentiating asset \u2014 now it must be activated and proven")
cards = [
 ("First in Ireland", "Virtue is the first private care home provider in Ireland to offer staff paid maternity benefits and health insurance \u2014 a genuine first-mover position in its category.", MAGENTA),
 ("A market under pressure", "Sector turnover is ~36%, the workforce is ~79% female, and a 50/50 non-EEA hiring rule limits international supply. Domestic attraction and retention are now critical.", BLUE),
 ("Costs are rising fast", "Monthly market spend on Indeed has risen ~4x since early 2024 and cost-per-click has nearly tripled \u2014 competing on spend alone is unsustainable.", ORANGE),
 ("Virtue already converts well", "Virtue's apply-start rate is 31% vs a 21% market average \u2014 the message resonates, but scale is small (5 live roles, 126th on spend).", TEAL),
 ("A clear benefit gap", "Only ~52% of Irish private employers offer maternity top-up, and ~20% in comparable low-wage sectors. Virtue is operating in open space.", PURPLE),
 ("The task is activation", "Make the benefit visible at every candidate touchpoint, prove it with employee stories, and measure the behavioural change before the market catches up.", DEEP),
]
cx, cy = M, 1.62
cw, chh, gx, gy = (CW-2*0.28)/3, 1.62, 0.28, 0.18
for i, (h, b, col) in enumerate(cards):
    r = i // 3; c = i % 3
    x = cx + c*(cw+gx); y = cy + r*(chh+gy)
    rect(s, x, y, cw, chh, fill=PALE)
    rect(s, x, y, 0.07, chh, fill=col)
    txt(s, x+0.18, y+0.14, cw-0.32, 0.4, [[(h, 12.5, col, True)]])
    txt(s, x+0.18, y+0.52, cw-0.32, chh-0.6, [[(b, 9.3, DARK, False)]], line_spacing=1.08)
footer(s, "Indeed market data (Jan 2024\u2013May 2026); BDO/NHI 2025; SOLAS; internal analysis", pg())
notes(s,
 "These are the seven points that matter. Virtue has moved first on a benefit that most of the market doesn't offer, in a sector that is structurally short of staff and getting more expensive to recruit in. The early signal is good - Virtue already converts interest into applications better than the market. But it's small, and the benefit isn't yet visible everywhere it needs to be.",
 "Push the room to treat this as a positioning opportunity with a closing window, not a perk that's been ticked off.",
 "Virtue is sitting on a differentiator most competitors can't yet match - the value is in activating and measuring it.",
 "If we achieved one thing with this benefit in the next 12 months, what would make it feel like a success to you?")

# ======================================================================
# Section 01 divider
# ======================================================================
divider("01", "The market context", "The Irish private care recruitment market is intensifying \u2014 more spend, higher cost, harder to stand out",
        "Let's start with the market Virtue is hiring into, because it shapes everything that follows. The short version: it's getting more crowded and more expensive.")

# ---- Slide: market spend & CPC ----
s = slide()
content_header(s, "Market context", "Competition for care candidates is intensifying \u2014 cost is climbing, not stabilising")
img_fit(s, os.path.join(CH, "A_market_spend_cpc.png"), M, 1.55, 6.2, 3.45, halign='left')
bx = 7.0
txt(s, bx, 1.62, 2.55, 0.5, [[("WHAT THE DATA SHOWS", 9, BLUE, True)]])
for i,(t1,t2) in enumerate([
   ("Market spend up ~4x","from ~\u20ac58k to ~\u20ac229k per month since Jan 2024"),
   ("Cost-per-click up ~3x","from \u20ac0.25 to \u20ac0.70 as more employers bid"),
   ("Same finite pool","~330\u2013360 employers competing for a flat candidate base")]):
    y=1.95+i*0.92
    txt(s, bx, y, 2.6, 0.3, [[(t1, 12, DARK, True)]])
    txt(s, bx, y+0.30, 2.6, 0.55, [[(t2, 9.5, GREY, False)]], line_spacing=1.08)
rect(s, bx, 4.62, 2.6, 0.5, fill=PALE)
txt(s, bx+0.14, 4.70, 2.35, 0.4, [[("Implication: ", 9.5, MAGENTA, True),("buying visibility alone is a losing race.", 9.5, DARK, False)]], line_spacing=1.05)
footer(s, "Indeed Ireland private care market, monthly, Jan 2024\u2013May 2026", pg())
notes(s,
 "This is the backdrop to everything else. Over the last two years the amount employers are spending in this market on Indeed has roughly quadrupled, and the cost of every click has nearly tripled. Crucially, the candidate pool hasn't grown to match - so everyone is paying more to reach the same people.",
 "This is why a benefit matters. When everyone is bidding up the same auction, the differentiator that wins is a better offer, not a bigger budget.",
 "You cannot out-spend this market indefinitely - you have to out-position it.",
 "How much has our own cost-per-hire moved over the last 18 months?")

# ---- Slide: candidate yield ----
s = slide()
content_header(s, "Market context", "Each role costs more but returns fewer applicants \u2014 efficiency is falling market-wide")
img_fit(s, os.path.join(CH, "B_candidate_yield.png"), M, 1.55, 6.2, 3.45, halign='left')
bx=7.0
txt(s, bx, 1.62, 2.6, 0.5, [[("THE SQUEEZE", 9, BLUE, True)]])
for i,(t1,t2) in enumerate([
   ("Investment per job up","from ~\u20ac29 to ~\u20ac82 per posted role"),
   ("Apply-starts per job down","candidate yield per role has compressed"),
   ("Quality beats quantity","standing out now drives efficiency, not budget")]):
    y=1.95+i*0.92
    txt(s, bx, y, 2.6, 0.3, [[(t1, 12, DARK, True)]])
    txt(s, bx, y+0.30, 2.6, 0.55, [[(t2, 9.5, GREY, False)]], line_spacing=1.08)
rect(s, bx, 4.62, 2.6, 0.5, fill=PALE)
txt(s, bx+0.14, 4.70, 2.35, 0.4, [[("Implication: ", 9.5, MAGENTA, True),("a stronger offer lifts conversion at no extra media cost.", 9.5, DARK, False)]], line_spacing=1.05)
footer(s, "Indeed Ireland private care market, monthly, Jan 2024\u2013May 2026", pg())
notes(s,
 "The same story from the employer's point of view. The cost of putting a role to market has risen sharply, while the number of people starting an application per role has drifted down. Employers are working harder for less. That's the definition of a market where differentiation, not budget, is the lever.",
 "Anything that lifts the conversion rate - a clearer, more compelling offer - effectively lowers cost-per-hire without spending another euro.",
 "The winning move is to convert better, and the benefit is exactly the kind of message that does that.",
 "Where in our own funnel do we lose the most candidates today?")

# ======================================================================
# Section 02 divider
# ======================================================================
divider("02", "Where Virtue stands", "Small in spend, strong in conversion \u2014 and getting more efficient",
        "Now let's place Virtue inside that market. The headline is a nice paradox: Virtue is tiny on spend but strong on the metrics that actually matter.")

# ---- Slide: competitive landscape ----
s = slide()
content_header(s, "Competitive landscape", "Virtue is a small spender in a market led by national home-care groups")
hbar_native(s, M, 1.62, 6.35, 3.35, C_ITEMS, namew=1.65, labw=0.62, vfmt=lambda v: f"\u20ac{v:.0f}k")
txt(s, M+1.65, 5.0, 4.2, 0.22, [[("Spend, Jan\u2013May 2026 (\u20ac000s) \u00b7 editable", 7.5, GREY, False)]])
bx=7.05
txt(s, bx, 1.62, 2.6, 0.5, [[("SHARE OF VOICE", 9, BLUE, True)]])
txt(s, bx, 1.95, 2.6, 1.3, [[("Comfort Keepers alone spends ", 10.5, DARK, False),("~140x", 10.5, BLUE, True),(" more than Virtue this year.", 10.5, DARK, False)]], line_spacing=1.12)
txt(s, bx, 3.05, 2.6, 1.2, [[("Virtue runs just ", 10.5, DARK, False),("5 live roles", 10.5, BLUE, True),(" and ranks 126th of 529 advertisers on spend.", 10.5, DARK, False)]], line_spacing=1.12)
rect(s, bx, 4.45, 2.6, 0.62, fill=PALE)
txt(s, bx+0.14, 4.53, 2.35, 0.5, [[("Opportunity: ", 9.5, MAGENTA, True),("low share of voice means targeted, benefit-led spend can move the needle quickly.", 9.0, DARK, False)]], line_spacing=1.04)
footer(s, "Indeed Ireland advertiser data, Jan\u2013May 2026 (n=529 advertisers)", pg())
notes(s,
 "Here's where Virtue sits among the people it's competing with for candidates. The market is led by large national home-care groups - Comfort Keepers, Danu, Dovida, Mowlam. Virtue is a very small spender by comparison, running a handful of roles. That's not a weakness in itself - it means there's a lot of headroom.",
 "Don't read low spend as low potential. It means Virtue is under-exposed, so well-targeted, benefit-led investment has room to deliver outsized gains.",
 "Virtue is punching at a fraction of its weight on visibility - there is clear room to grow share of voice.",
 "Which two or three locations or roles would you most want to win if we increased visibility there?")

# ---- Slide: apply-start rate ----
s = slide()
content_header(s, "Conversion performance", "When candidates see Virtue, they act \u2014 conversion runs well above the market")
col_native(s, M+0.1, 1.72, 6.15, 2.82, D_ITEMS, vfmt=lambda v: f"{v:.1f}%")
txt(s, M, 5.02, 4.5, 0.22, [[("Apply-start rate (%) \u00b7 Jan\u2013May 2026 \u00b7 editable", 7.5, GREY, False)]])
bx=7.05
txt(s, bx, 1.62, 2.6, 0.5, [[("THE STANDOUT METRIC", 9, BLUE, True)]])
txt(s, bx, 1.98, 2.6, 1.1, [[("31%", 30, BLUE, True)]])
txt(s, bx, 2.62, 2.6, 0.8, [[("apply-start rate vs ", 10.5, DARK, False),("21%", 10.5, DARK, True),(" market average \u2014 ~47% higher.", 10.5, DARK, False)]], line_spacing=1.12)
rect(s, bx, 3.95, 2.6, 1.1, fill=PALE)
txt(s, bx+0.14, 4.05, 2.35, 0.95, [[("Read-across: ", 9.5, MAGENTA, True),("the offer already converts. The constraint is reach and consistency, not appeal. More visibility on a message that works should compound.", 9.0, DARK, False)]], line_spacing=1.05)
footer(s, "Indeed Ireland advertiser data, Jan\u2013May 2026", pg())
notes(s,
 "This is the most encouraging slide in the pack. Of the people who click a Virtue ad, 31% go on to start an application - against a market average of 21%. That's roughly half as much again. It tells us the proposition is landing: when candidates encounter Virtue, they're more likely to act than they are with competitors.",
 "This reframes the whole conversation. Virtue's problem isn't appeal, it's reach. We have a message that converts - we just need more people to see it, more often, more consistently.",
 "Virtue already converts interest into action better than the market - scale the reach and the results should follow.",
 "Do we know what candidates say is the reason they apply to Virtue rather than a competitor?")

# ---- Slide: efficiency ----
s = slide()
content_header(s, "Spend efficiency", "Virtue is buying candidates more cheaply than it did six months ago")
img_fit(s, os.path.join(CH, "E_virtue_efficiency.png"), M, 1.55, 6.1, 3.5, halign='left')
bx=6.9
txt(s, bx, 1.62, 2.7, 0.5, [[("MOVING THE RIGHT WAY", 9, BLUE, True)]])
for i,(t1,t2) in enumerate([
   ("Cost-per-click \u20ac1.82 \u2192 \u20ac0.83","more than halved"),
   ("Cost per apply-start \u20ac5.68 \u2192 \u20ac2.68","down ~53%"),
   ("Conversion held at ~31%","efficiency gained without losing quality")]):
    y=1.95+i*0.95
    txt(s, bx, y, 2.75, 0.35, [[(t1, 11.5, DARK, True)]])
    txt(s, bx, y+0.30, 2.75, 0.4, [[(t2, 9.5, GREY, False)]])
rect(s, bx, 4.7, 2.75, 0.45, fill=PALE)
txt(s, bx+0.14, 4.77, 2.5, 0.35, [[("A strong base to scale from.", 9.5, MAGENTA, True)]])
footer(s, "Indeed Ireland advertiser data \u2014 Virtue Integrated Care, H2 2025 vs Jan\u2013May 2026", pg())
notes(s,
 "Comparing the second half of last year with this year, Virtue is getting more for its money. Cost-per-click has more than halved, the cost of generating an application has dropped by about half, and the conversion rate has held up at around 31%. So this isn't cheaper traffic that converts worse - it's genuine efficiency.",
 "This matters for the budget conversation: Virtue has shown it can spend smartly. That de-risks the case for investing a bit more behind the benefit message.",
 "Virtue's media is already efficient and improving - that's a credible platform to scale investment from.",
 "If the numbers continue to improve, what would justify increasing the employer-brand budget?")

# ======================================================================
# Section 03 divider
# ======================================================================
divider("03", "Candidate demand by location", "Where the candidates are \u2014 and where supply is tightest",
        "If we're going to invest behind the benefit, we should know where to point it. This section maps where the candidates are and where they're hardest to find.")

# ---- Slide: county demand (all 26 counties) ----
s = slide()
content_header(s, "Location opportunity", "All 26 counties: demand clusters in Dublin & Cork \u2014 supply is tightest in the west & border", headline_size=17)
txt(s, M, 1.40, CW, 0.3,
    [[("Scale plays in the deep eastern pools; lead with the benefit where seekers-per-job are lowest (magenta) \u2014 Donegal, Clare, Sligo, Cavan, Roscommon, Cork.", 9.5, GREY, False)]],
    line_spacing=1.05)
img_fit(s, os.path.join(CH, "F_county_demand.png"), M, 1.78, CW, 3.4, halign='center', valign='top')
footer(s, "Indeed Ireland \u2014 care-role jobseekers & jobs by county (26 counties; 'unknown' excluded)", pg())
notes(s,
 "This is every county, so nothing is hidden. Read it two ways. First, volume - the bar length and the first number show where the candidates are: Dublin is in a league of its own, then Cork, Kildare, Meath and the eastern commuter belt. Those are your scale markets. Second, scarcity - the colour and the 'per job' figure show how tight supply is. The magenta counties have the fewest jobseekers per advertised role: Donegal, Clare, Sligo, Cavan and Roscommon in the west and border, plus Cork despite its size. Those are the hardest to fill, and that's exactly where a differentiated offer like the benefit does the most work.",
 "Match the tactic to the geography. Use reach and efficiency in the deep eastern pools; use the benefit as the wedge in the tight-supply counties where everyone is fishing in a small pond and a bigger bid alone won't win.",
 "Volume tells you where to scale; scarcity tells you where the benefit matters most - and the two don't always overlap.",
 "Which of these counties cause us the most recruitment pain today, and are we investing there in proportion to the difficulty?")

# ======================================================================
# Section 04 divider
# ======================================================================
divider("04", "The benefit & the EVP", "What the maternity and health benefit means to candidates, the EVP it proves, and where the UK shows it's heading",
        "This is the core of the case: what the maternity and health benefit actually means to the people Virtue wants to hire, how it anchors the wider employer brand, and why now is the moment to lean into it.")

# ---- Slide: benefit gap ----
s = slide()
content_header(s, "The benefit gap", "Virtue is operating in open space the rest of the market has left empty")
img_fit(s, os.path.join(CH, "G_benefit_gap.png"), M, 1.55, 6.0, 3.5, halign='left')
bx=6.8
txt(s, bx, 1.62, 2.75, 0.5, [[("WHY IT MATTERS", 9, BLUE, True)]])
txt(s, bx, 1.96, 2.8, 1.2, [[("Only ~52% of Irish private employers offer maternity top-up \u2014 and ~20% in comparable low-wage sectors.", 10.5, DARK, False)]], line_spacing=1.14)
txt(s, bx, 3.15, 2.8, 1.2, [[("With ~79% of the care workforce female and mostly mid-career, this benefit speaks directly to the majority of Virtue's candidates.", 10.5, DARK, False)]], line_spacing=1.14)
rect(s, bx, 4.5, 2.8, 0.6, fill=PALE)
txt(s, bx+0.14, 4.58, 2.55, 0.5, [[("This is category-level differentiation, not an incremental perk.", 9.3, MAGENTA, True)]], line_spacing=1.05)
footer(s, "cleareye.ie maternity benefit benchmarks; SOLAS care workforce profile; internal analysis", pg())
notes(s,
 "This is the commercial heart of the case. Across the Irish private sector only about half of employers offer any maternity top-up, and in low-wage sectors comparable to care it's around one in five. Virtue offering it - alongside health insurance - puts it in space the rest of the market has largely left empty. And because nearly four in five care workers are women, mostly mid-career, this isn't a niche perk - it's relevant to the majority of the people Virtue wants to hire.",
 "Position this as a category move. Most competitors are still competing on near-identical pay and 'supportive team' messaging - Virtue is saying something structurally different.",
 "Virtue is offering something most of the market does not, to exactly the audience that values it most.",
 "Have we confirmed the precise terms - weeks, percentage of salary, eligibility - so we can state them with confidence?")

# ---- Slide: first-in-Ireland recruitment & retention dividend ----
s = slide()
content_header(s, "First-mover advantage", "First in Irish private care \u2014 a measurable recruitment and retention dividend")
# impact chain across the top
chain = [
 ("First to offer it", "A claim no Irish private care competitor can currently match", MAGENTA),
 ("Stand out to scarce candidates", "A genuine differentiator where supply is tightest", PURPLE),
 ("More & stronger applicants", "Already converting at 31% vs 21% market", BLUE),
 ("Higher offer acceptance", "Total-reward advantage at the decision point", TEAL),
 ("Lower early attrition", "Loyalty built before day one \u2014 fewer replacements", DEEP),
]
cwm = (CW - 4*0.16) / 5
for i, (h, b, col) in enumerate(chain):
    x = M + i*(cwm+0.16); y = 1.62
    rect(s, x, y, cwm, 1.5, fill=PALE)
    rect(s, x, y, cwm, 0.1, fill=col)
    txt(s, x+0.12, y+0.22, cwm-0.24, 0.6, [[(h, 10, col, True)]], line_spacing=0.98)
    txt(s, x+0.12, y+0.86, cwm-0.24, 0.6, [[(b, 7.8, GREY, False)]], line_spacing=1.05)
    if i < len(chain)-1:
        txt(s, x+cwm-0.02, y+0.55, 0.18, 0.4, [[("\u203a", 16, RGBColor(0xB9,0xC1,0xD2), True)]])
# perspective panels
py = 3.42
txt(s, M, py, 4.3, 0.3, [[("WHY IT MATTERS NOW", 9, BLUE, True)]])
for i, t in enumerate([
   "Sector turnover ~36% and rising \u2014 every retained hire avoids a costly replacement",
   "~79% of the care workforce is female and mid-career \u2014 the benefit speaks to the majority",
   "A 50/50 non-EEA hiring cap limits international supply \u2014 domestic retention is now decisive"]):
    yy = py+0.34+i*0.42
    rect(s, M, yy+0.05, 0.07, 0.07, fill=BLUE, shape=MSO_SHAPE.OVAL)
    txt(s, M+0.2, yy-0.02, 4.1, 0.42, [[(t, 8.9, DARK, False)]], line_spacing=1.02)
ox = 5.05
txt(s, ox, py, 4.4, 0.3, [[("THE COMMERCIAL PAYOFF", 9, MAGENTA, True)]])
for i, t in enumerate([
   "Replacing a care worker costs ~1\u20133 months' salary \u2014 small retention gains compound fast",
   "A 5% lift in acceptance and a 5% cut in early attrition together lower true cost-per-hire",
   "First-mover window is time-limited \u2014 the advantage erodes once competitors copy it"]):
    yy = py+0.34+i*0.42
    rect(s, ox, yy+0.05, 0.07, 0.07, fill=MAGENTA, shape=MSO_SHAPE.OVAL)
    txt(s, ox+0.2, yy-0.02, 4.25, 0.42, [[(t, 8.9, DARK, False)]], line_spacing=1.02)
footer(s, "BDO/NHI 2025 (turnover); SOLAS (workforce); Indeed conversion data; internal analysis", pg())
notes(s,
 "Let's translate 'first in Ireland' into what it actually does for hiring. Being first gives Virtue a claim no competitor can currently match, and that matters most exactly where candidates are scarcest. The chain along the top is the logic: a differentiated offer makes Virtue stand out, which brings more and stronger applicants - and we can already see that in the 31% apply-start rate. That feeds higher offer acceptance because candidates perceive a better total reward, and then lower early attrition because people feel invested in before they even start. The two panels below are the market perspective: why it matters now - turnover near 36%, a mostly female mid-career workforce, and a regulatory cap on international hiring that makes domestic retention decisive - and the commercial payoff, where replacement costs of one to three months' salary mean even small retention gains pay back quickly.",
 "Frame 'first' as a recruitment and retention asset with a shelf-life, not a bragging right. The value is the head-start, and it decays the moment a competitor matches it.",
 "Being first turns the benefit into both an attraction magnet and a retention anchor - but the clock is running.",
 "If we treated this as our single biggest retention lever this year, what would we measure to prove it's working?")

# ---- Slide: why candidates care ----
s = slide()
content_header(s, "What it means to candidates", "Candidates don't read benefits \u2014 they scan for signals, and this one lands")
lx = M
txt(s, lx, 1.62, 4.1, 0.4, [[("HOW CANDIDATES RESPOND", 9, BLUE, True)]])
pts = [
 ("A scroll-stop signal", "A maternity line in the opening paragraph says \u2018this employer thinks about me beyond my shift\u2019 \u2014 in a feed of near-identical care ads."),
 ("It removes real anxiety", "For a carer in her late 20s\u201330s, salary continuity during maternity addresses one of the most stressful financial moments in adult life."),
 ("Credibility beats cost", "The emotional credibility this creates is disproportionate to the actual cost of the benefit \u2014 and it is mainstream, not niche."),
]
for i, (h, b) in enumerate(pts):
    y = 2.0 + i*1.02
    txt(s, lx, y, 4.0, 0.3, [[(h, 12.5, DARK, True)]])
    txt(s, lx, y+0.30, 4.0, 0.65, [[(b, 9.5, GREY, False)]], line_spacing=1.1)
rx = 4.95
txt(s, rx, 1.62, 4.45, 0.4, [[("WHO RESPONDS MOST", 9, BLUE, True)]])
segs = [
 ("Healthcare Assistants (female, 25\u201340)", "Primary audience \u2014 planning or have planned families"),
 ("Return-to-work candidates", "Need proof the family-friendly culture is genuine"),
 ("Candidates in low-benefit roles", "Comparison-shopping; currently receive no top-up"),
 ("International healthcare workers", "Want a stable employer with formal benefit structures"),
 ("Movers from retail / hospitality", "Benefit compares very favourably to their current job"),
]
for i, (h, b) in enumerate(segs):
    y = 1.98 + i*0.625
    rect(s, rx, y, 4.45, 0.55, fill=PALE)
    rect(s, rx, y, 0.06, 0.55, fill=A11Y)
    txt(s, rx+0.18, y+0.07, 4.2, 0.25, [[(h, 9.8, DEEP, True)]])
    txt(s, rx+0.18, y+0.30, 4.2, 0.22, [[(b, 8.5, GREY, False)]])
footer(s, "SOLAS care workforce profile (~79% female); candidate-behaviour analysis; internal segmentation", pg())
notes(s,
 "Two things leadership should take from this. First, how candidates behave: people don't read a benefits list, they scan for a signal. A maternity line near the top of an ad is a scroll-stop in a feed where every ad looks the same, and for the typical candidate - a woman in her late twenties or thirties - it speaks to a genuine financial worry. The emotional pull is far bigger than the cost. Second, who responds: this isn't niche. From front-line HCAs to returners, to people coming from retail and hospitality, the benefit is relevant to the majority of Virtue's pipeline.",
 "Reframe the benefit from 'a perk we offer' to 'a signal that changes who applies'. It speaks to the dominant candidate profile, not a minority.",
 "This benefit is mainstream to Virtue's audience - it should be front and centre, not a line at the bottom of the ad.",
 "Do our current ads put the benefit where candidates actually look - the first paragraph - or is it buried in a list?")

# ---- Slide: EVP pillars ----
s = slide()
content_header(s, "Employer value proposition", "The benefit is proof of a bigger promise \u2014 five EVP pillars, one consistent story")
pillars = [
 ("We care for you", "Maternity top-up, health insurance, wellbeing support", MAGENTA, True),
 ("We're with you long-term", "Career development, progression and stability", BLUE, False),
 ("We respect your life", "Family-friendly, flexible, supportive culture", TEAL, False),
 ("We're proud of our work", "Quality, values-led care for residents", PURPLE, False),
 ("You belong here", "Inclusive, team-oriented and welcoming", ORANGE, False),
]
pw = (CW - 4*0.2) / 5
for i, (h, b, col, proof) in enumerate(pillars):
    x = M + i*(pw+0.2); y = 1.66
    rect(s, x, y, pw, 2.18, fill=PALE)
    rect(s, x, y, pw, 0.12, fill=col)
    txt(s, x+0.13, y+0.30, pw-0.26, 0.85, [[(h, 11.5, col, True)]], line_spacing=1.0)
    txt(s, x+0.13, y+1.15, pw-0.26, 0.85, [[(b, 8.8, DARK, False)]], line_spacing=1.1)
    if proof:
        txt(s, x+0.13, y+1.86, pw-0.26, 0.28, [[("\u25b2 the benefit lives here", 7.3, MAGENTA, True)]])
rect(s, M, 4.12, CW, 0.95, fill=NAVY)
txt(s, M+0.28, 4.2, CW-0.56, 0.8,
    [[("\u201cAt Virtue, caring is not just what we do for our residents. It is what we do for each other \u2014 we invest in your health, your family and your future, because when you're supported, the people in our care are too.\u201d", 11.5, WHITE, False)]],
    line_spacing=1.12, anchor=MSO_ANCHOR.MIDDLE)
footer(s, "Virtue EVP framework \u2014 internal analysis", pg())
notes(s,
 "The benefit is most powerful when it's the proof point for a bigger promise, not a stand-alone perk. We'd anchor Virtue's employer brand on five pillars - we care for you, we're with you long-term, we respect your life outside work, we're proud of our work, and you belong here. The maternity and health benefit is the tangible proof of the first pillar, and it makes all five more credible. The line along the bottom is the core candidate promise we'd lead with everywhere.",
 "The message to land is consistency: every channel, from a job ad to the offer letter, should tell one version of this story, with the benefit as proof rather than a bullet.",
 "Lead with the human promise and use the benefit as the proof - 'we invest in your health, your family and your future'.",
 "Does our messaging today tell one consistent story, or does each channel say something slightly different?")

# ---- Slide: UK benchmark table ----
s = slide()
content_header(s, "UK benchmark", "The UK shows where Irish private care is heading \u2014 Virtue can be there first")
rows = [
 ("Employer", "Maternity / family benefit", "Health & wellbeing"),
 ("Barchester Healthcare", "Enhanced maternity & paternity pay", "Listed benefits package"),
 ("HC-One", "Enhanced maternity, paternity & adoption leave (2023)", "Free 24/7 GP access, EAP"),
 ("Nuffield Health", "10 weeks full pay, 16 weeks half pay + SMP", "Free healthcare plan & gym"),
 ("Bupa", "Up to 52 weeks; 13 wks basic pay, 8\u201313 wks half pay", "Employer-funded health insurance"),
]
tx, ty = M, 1.62; tw = 6.4
colw = [1.55, 3.05, 1.8]
rh = 0.62
for ri, row in enumerate(rows):
    y = ty + ri*rh
    if ri == 0:
        rect(s, tx, y, tw, rh, fill=BLUE)
    elif ri % 2 == 0:
        rect(s, tx, y, tw, rh, fill=PALE)
    cx = tx
    for ci, cell in enumerate(row):
        col = WHITE if ri == 0 else DARK
        bold = ri == 0 or ci == 0
        txt(s, cx+0.12, y+0.06, colw[ci]-0.18, rh-0.1,
            [[(cell, 9.0 if ri else 9.8, (BLUE if (ci==0 and ri) else col), bold)]],
            anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        cx += colw[ci]
rect(s, tx, ty, tw, rh*len(rows), line=LGREY, line_w=0.75)
bx=7.15
txt(s, bx, 1.62, 2.5, 0.5, [[("THE SIGNAL", 9, BLUE, True)]])
txt(s, bx, 1.96, 2.5, 1.6, [[("UK private care has used enhanced family and health benefits as employer-brand tools for years \u2014 bundled into a 'we invest in the whole person' story.", 10, DARK, False)]], line_spacing=1.14)
rect(s, bx, 3.75, 2.5, 1.3, fill=PALE)
txt(s, bx+0.14, 3.85, 2.25, 1.15, [[("First-mover window: ", 9.5, MAGENTA, True),("Ireland's market is likely 3\u20135 years behind. Acting now \u2014 before benefits normalise \u2014 is a time-limited advantage.", 9.0, DARK, False)]], line_spacing=1.08)
footer(s, "Barchester, HC-One, Nuffield Health, Bupa published benefits (2022\u20132023)", pg())
notes(s,
 "To sense-check whether this is the right bet, look at the UK. The leading private care providers there - Barchester, HC-One, Nuffield Health, Bupa - have offered enhanced maternity and health benefits as employer-brand tools for years, and they bundle them into a bigger 'we invest in the whole person' story. The UK is typically three to five years ahead of Ireland on this. So Virtue isn't taking an untested gamble - it's getting to where the market is heading first.",
 "Use the UK as proof, not as a copy exercise. The lesson is to bundle the benefit into a values story, and to move before the Irish market normalises it.",
 "What's standard in the UK today is where Ireland is going - Virtue can own that position before competitors arrive.",
 "Which Irish competitor do you think is most likely to copy this, and how fast?")

# ======================================================================
# Section 05 divider \u2014 Getting it to candidates
# ======================================================================
divider("05", "Getting it to candidates", "Make the benefit visible and consistent at every step \u2014 from first impression to offer",
        "Knowing the benefit works isn't enough - it has to reach candidates consistently. This section is the practical activation: where the message shows up, on which Indeed surface, and what we measure.")

# ---- Slide: candidate journey map ----
s = slide()
content_header(s, "Candidate journey", "One message, every touchpoint \u2014 mapped to the Indeed surface that carries it")
jrows = [
 ("Stage", "What Virtue says", "Indeed / Glassdoor surface", "Success metric"),
 ("Awareness", "\u2018Benefits most Irish care employers don't offer\u2019", "Employer Branded Ads, social", "Brand search lift"),
 ("Job search", "Benefit in the headline & first paragraph", "Sponsored & Premium Jobs", "Click-through, apply rate"),
 ("Consideration", "Employee stories + EBH proof points", "EBH, Company Page, Glassdoor", "Time on page"),
 ("Application", "Reinforce benefit; frictionless apply", "Indeed Apply", "Completion rate"),
 ("Interview", "Recruiter talk track: specific terms & stories", "Interview, Hiring Events", "Interview-to-offer"),
 ("Offer", "Restate the benefit as total reward", "Offer letter, recruiter call", "Offer acceptance"),
 ("Onboard & advocacy", "Reinforce; encourage honest reviews", "Onboarding, Glassdoor", "90-day retention, reviews"),
]
tx, ty = M, 1.5; tw = CW
jcolw = [1.55, 3.15, 2.4, 1.8]
rh = 0.375
for ri, row in enumerate(jrows):
    y = ty + ri*rh
    if ri == 0:
        rect(s, tx, y, tw, rh, fill=BLUE)
    elif ri % 2 == 0:
        rect(s, tx, y, tw, rh, fill=PALE)
    cx = tx
    for ci, cell in enumerate(row):
        colr = WHITE if ri == 0 else DARK
        bold = ri == 0 or ci == 0
        txt(s, cx+0.12, y+0.03, jcolw[ci]-0.18, rh-0.04,
            [[(cell, 8.6 if ri == 0 else 8.7, (BLUE if (ci == 0 and ri) else colr), bold)]],
            anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.98)
        cx += jcolw[ci]
rect(s, tx, ty, tw, rh*len(jrows), line=LGREY, line_w=0.75)
by = ty + rh*len(jrows) + 0.09
rect(s, tx, by, tw, 0.4, fill=NAVY)
txt(s, tx+0.16, by+0.03, tw-0.3, 0.34,
    [[("Recommended opening line:  ", 8.8, A11Y, True),
      ("\u201cVirtue offers paid maternity leave top-up and health insurance \u2014 because caring for our team matters as much as caring for our residents.\u201d", 8.8, WHITE, False)]],
    anchor=MSO_ANCHOR.MIDDLE)
footer(s, "Candidate journey messaging map across Indeed & Glassdoor surfaces \u2014 internal analysis", pg())
notes(s,
 "This is the how-to. Reading across the journey: in awareness we use brand ads to plant the message; in job search the benefit sits in the headline and first paragraph of Sponsored and Premium jobs; in consideration the EBH and employee stories prove it's real; through application, interview and offer we keep restating it - including in the offer letter as part of total reward; and after hire we reinforce it and invite honest reviews. Each stage has an owner surface and a metric, so we can see where it's working and where it leaks. The line at the bottom is the recommended opening line for ads.",
 "The single biggest failure mode is inconsistency - the benefit in the ad but missing from the interview or offer. This map is the standard that prevents that.",
 "Consistency is the whole game: the same proof, on every surface, measured at every stage.",
 "If we audited our candidate journey today, where would the benefit message disappear?")

# ---- Slide: marginal gains ----
s = slide()
content_header(s, "The commercial case", "Small gains at each funnel stage compound into a material return")
img_fit(s, os.path.join(CH, "H_marginal_gains.png"), M, 1.55, 6.2, 3.5, halign='left')
bx=6.95
txt(s, bx, 1.62, 2.65, 0.5, [[("MARGINAL GAINS", 9, BLUE, True)]])
txt(s, bx, 1.96, 2.7, 1.5, [[("A 5% relative lift at each stage \u2014 from click to offer acceptance \u2014 compounds into more hires from the same or lower spend.", 10.5, DARK, False)]], line_spacing=1.14)
txt(s, bx, 3.45, 2.7, 1.0, [[("Replacing a care worker costs ~1\u20133 months' salary, so even small retention gains pay back quickly.", 10.5, DARK, False)]], line_spacing=1.14)
rect(s, bx, 4.62, 2.7, 0.5, fill=PALE)
txt(s, bx+0.14, 4.70, 2.45, 0.4, [[("Illustrative \u2014 apply Virtue's actual volumes.", 8.8, MAGENTA, True)]])
footer(s, "Illustrative marginal-gains model; replacement-cost estimate per internal analysis", pg())
notes(s,
 "This is how a benefit becomes a number the board cares about. If the benefit lifts each stage of the funnel by just 5% - click-through, apply-start, completion, interview attendance, offer acceptance - those small gains multiply. You end up with more hires from the same or lower spend. And because replacing a care worker costs one to three months' salary, even a small improvement in early retention pays for itself quickly.",
 "These figures are illustrative - the point is the mechanism, not the exact numbers. We should plug in Virtue's real volumes to size the prize.",
 "You don't need a big effect at any one stage - small, compounding gains across the funnel are what move cost-per-hire.",
 "What's the single funnel metric you'd most want to move first - applications, acceptance, or early retention?")

# ======================================================================
# Section 06 divider
# ======================================================================
divider("06", "Risks & recommended actions", "Protect the advantage, then activate it",
        "Finally, what could go wrong, and what we'd actually do about it - a clear, sequenced plan the room can say yes to.")

# ---- Slide: risks & opportunities ----
s = slide()
content_header(s, "Risks & opportunities", "The benefit is real \u2014 the risk is that it stays a bullet point")
colw2 = (CW-0.4)/2
# Risks
rx = M
rect(s, rx, 1.6, colw2, 0.5, fill=MAGENTA)
txt(s, rx+0.16, 1.68, colw2-0.3, 0.35, [[("RISKS / WATCH-OUTS", 11, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
risks = [
 "Message visible in ads but absent from EBH, interview and offer \u2014 diluting the investment",
 "Over-claiming 'first in Ireland' without verified competitor checks",
 "Candidates not understanding the value (terms, eligibility) without a clear FAQ",
 "Treated as an HR update, not a commercial, C-suite-owned asset",
 "No baseline or measurement \u2014 leaving it exposed at the next budget cycle",
]
for i,r in enumerate(risks):
    y=2.2+i*0.6
    rect(s, rx, y+0.06, 0.08, 0.08, fill=MAGENTA, shape=MSO_SHAPE.OVAL)
    txt(s, rx+0.2, y-0.02, colw2-0.25, 0.55, [[(r, 9.3, DARK, False)]], line_spacing=1.04)
# Opportunities
ox = M+colw2+0.4
rect(s, ox, 1.6, colw2, 0.5, fill=BLUE)
txt(s, ox+0.16, 1.68, colw2-0.3, 0.35, [[("OPPORTUNITIES", 11, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
opps = [
 "Move the benefit to paragraph one of every ad \u2014 it already converts at 31%",
 "Build the EBH into a proof destination with real employee stories",
 "Lead with the benefit in tight-supply counties (Cork, Donegal, Clare, Sligo)",
 "Scale efficient, benefit-led spend \u2014 Virtue is only 126th on share of voice",
 "Set a before/after measurement model to prove ROI to the board",
]
for i,r in enumerate(opps):
    y=2.2+i*0.6
    rect(s, ox, y+0.06, 0.08, 0.08, fill=BLUE, shape=MSO_SHAPE.OVAL)
    txt(s, ox+0.2, y-0.02, colw2-0.25, 0.55, [[(r, 9.3, DARK, False)]], line_spacing=1.04)
footer(s, "Internal analysis; Indeed market data", pg())
notes(s,
 "Quick balance sheet. The risks are mostly about inconsistency: if the benefit appears in the ad but not in the EBH, the interview or the offer, it loses credibility and the investment leaks away. There's also a reputational risk in over-claiming 'first in Ireland', and a commercial risk in not measuring it. The opportunities mirror those: make it visible everywhere, prove it with real employee voices, target the tight-supply counties, scale the spend behind a message that already works, and measure it.",
 "The biggest risk isn't the benefit failing - it's the benefit being under-used. Treat consistency and measurement as the priorities.",
 "A great benefit that candidates only meet once is a bullet point - the value comes from making it impossible to miss.",
 "Who in the business should own this end-to-end, so it doesn't fall between HR, marketing and recruitment?")

# ---- Slide: 30/60/90 ----
s = slide()
content_header(s, "Recommended actions", "A 90-day plan: optimise what's live, expand across the journey, then prove it")
phases = [
 ("DAYS 1\u201330", "Optimise what's live", BLUE, [
   "Move the benefit to paragraph one of every live ad; A/B test vs a control",
   "Update the EBH and Company Page with specific benefit terms",
   "Brief every recruiter and hiring manager on the talk track",
   "Set baseline metrics: CTR, apply-start, completion, offer acceptance, 90-day retention",
 ]),
 ("DAYS 31\u201360", "Expand across the journey", TEAL, [
   "Produce 2\u20133 employee stories; publish on EBH, social and Glassdoor",
   "Launch benefit-led brand ads to passive candidates in priority counties",
   "Activate Premium Sponsored Jobs for the hardest-to-fill roles",
   "Add the benefit to offer letters and onboarding",
 ]),
 ("DAYS 61\u201390", "Prove impact & scale", MAGENTA, [
   "Review before/after data; present the marginal gains to leadership",
   "Run a Hiring Event in a tight-supply county",
   "Scale the best-performing copy across all roles",
   "Prepare the 90-day ROI summary and next-investment case for the board",
 ]),
]
pw = (CW-2*0.3)/3
for i,(tag,title,col,items) in enumerate(phases):
    x = M + i*(pw+0.3)
    rect(s, x, 1.6, pw, 0.72, fill=col)
    txt(s, x+0.16, 1.66, pw-0.3, 0.3, [[(tag, 10.5, WHITE, True)]])
    txt(s, x+0.16, 1.95, pw-0.3, 0.34, [[(title, 12.5, WHITE, True)]])
    rect(s, x, 2.32, pw, 2.62, fill=PALE)
    for j,it in enumerate(items):
        y=2.46+j*0.62
        rect(s, x+0.16, y+0.05, 0.07, 0.07, fill=col, shape=MSO_SHAPE.OVAL)
        txt(s, x+0.34, y-0.02, pw-0.5, 0.6, [[(it, 8.8, DARK, False)]], line_spacing=1.05)
footer(s, "Recommended activation roadmap \u2014 internal analysis", pg())
notes(s,
 "The plan is deliberately sequenced: optimise, expand, prove. In the first month we fix what's already live - get the benefit to the top of every ad, update the EBH, brief recruiters, and set a baseline so we can measure change. In the second month we expand across the journey with employee stories, brand ads in priority counties, premium placements for hard roles, and updated offer letters. In the final month we prove it - review the before-and-after data, run a hiring event in a tough county, scale what works, and take an ROI case to the board.",
 "Note the order. We're not asking for a big upfront bet - we optimise cheaply first, show signal, then scale with evidence. That's a low-risk path to a bigger commitment.",
 "Optimise first, expand second, prove third - each phase de-risks the next.",
 "Of these actions, which can we commit to starting in the next 30 days?")

# ======================================================================
# Closing
# ======================================================================
s = slide()
rect(s, -0.1, -0.1, 10.2, 5.825, fill=NAVY)
rect(s, 0, 0, 0.16, 5.625, fill=A11Y)
img(s, os.path.join(AS, "indeed_white.png"), M, 0.55, w=0.9)
txt(s, M, 2.05, 8.8, 1.4, [[("The benefit is the investment.", 27, WHITE, True)],
                            [("The candidate journey is where Virtue proves the return.", 27, A11Y, True)]], line_spacing=1.05)
txt(s, M, 3.75, 8.6, 0.9, [[("Activate it consistently, prove it with data, and turn it into a market position before the market catches up.", 13.5, RGBColor(0xD7,0xE0,0xF2), False)]], line_spacing=1.2)
txt(s, M, 5.05, 8.8, 0.3, [[("Virtue Integrated Care Services  ·  Market Overview & Talent Strategy  ·  Confidential", 9.5, RGBColor(0x9F,0xB2,0xD6), False)]])
notes(s,
 "To close: Virtue hasn't just launched a benefit, it's created an employer-brand asset in a market that's getting harder and more expensive to hire in. The data shows the offer already converts above the market. The job now is to make it visible everywhere, prove it with employee stories and before-and-after measurement, and turn it into a position competitors can't easily copy.",
 "Leave the room with a decision, not just agreement - who owns it, what we measure, and what we commit to in the next 30 days.",
 "Market positions, unlike benefits, are hard to copy - the window is open but it won't stay open forever.")

prs.save(OUT)
print("SAVED:", OUT)
print("Slides:", len(prs.slides._sldIdLst))
